"""
OmniMentor — CS 6460 Final Presentation Generator v2
15 slides, GT Development Track (hourglass structure).
Visuals: TPM role diagram, system block diagram, 4-step flowchart,
         ablation bar chart, domain corpus cards.
"""
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ─────────────────────────────────────────────────────────────────
BG      = RGBColor(0x06, 0x0B, 0x18)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BLUE    = RGBColor(0x3B, 0x82, 0xF6)
PURPLE  = RGBColor(0x8B, 0x5C, 0xF6)
GREEN   = RGBColor(0x10, 0xB9, 0x81)
AMBER   = RGBColor(0xF5, 0x9E, 0x0B)
DIM     = RGBColor(0xA0, 0xAE, 0xC0)
CARD    = RGBColor(0x11, 0x18, 0x27)
CARD2   = RGBColor(0x0D, 0x14, 0x22)
DARK2   = RGBColor(0x0A, 0x12, 0x24)

# ── Low-level helpers ────────────────────────────────────────────────────────
def _bg(slide):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = BG

def _tb(slide, l, t, w, h, text, sz=15, color=WHITE, bold=False,
        italic=False, align=PP_ALIGN.LEFT, wrap=True):
    s = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = s.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    return s

def _rect(slide, l, t, w, h, fill, border=None, label="",
          lsz=12, lcol=WHITE, lbold=False, wrap=True):
    s = slide.shapes.add_shape(
        1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(1.2)
    else:
        s.line.fill.background()
    if label:
        tf = s.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.text = label
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(lsz)
        p.font.color.rgb = lcol
        p.font.bold = lbold
    return s

def _oval(slide, l, t, w, h, fill, border=WHITE, label="", lsz=12, lcol=WHITE):
    s = slide.shapes.add_shape(
        9, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = border
    s.line.width = Pt(1)
    if label:
        tf = s.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(lsz)
        p.font.color.rgb = lcol
        p.font.bold = True
    return s

def _arrow_r(slide, l, t, w=0.32, h=0.22):
    s = slide.shapes.add_shape(13, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = BLUE; s.line.fill.background()

def _arrow_d(slide, l, t, w=0.22, h=0.32):
    s = slide.shapes.add_shape(36, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = BLUE; s.line.fill.background()

def _title(slide, text, sz=33):
    _tb(slide, 0.4, 0.18, 9.2, 0.72, text, sz=sz, color=WHITE, bold=True)
    bar = slide.shapes.add_shape(1, Inches(0.4), Inches(0.88), Inches(9.2), Inches(0.04))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()

def _footer(slide):
    _tb(slide, 0.4, 7.18, 9.2, 0.32,
        "Arvind Sharma  ·  CS 6460  ·  Georgia Institute of Technology  ·  Spring 2026",
        sz=11, color=BLUE, align=PP_ALIGN.CENTER)

def _callout(slide, l, t, w, h, text, sz=14, border=BLUE):
    s = _rect(slide, l, t, w, h, DARK2, border=border)
    tf = s.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(sz)
    p.font.color.rgb = WHITE

def _hbar(slide, l, t, total_w, filled_w, h, fill):
    _rect(slide, l, t, total_w, h, RGBColor(0x1A, 0x24, 0x3A))
    _rect(slide, l, t, filled_w, h, fill)


def _set_notes(slide, text):
    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    notes.text = text


SPEAKER_NOTES = [
    {
        "title": "OmniMentor",
        "notes": (
            "Open with one sentence: OmniMentor is a guided practice environment for new TPMs "
            "to build architecture reasoning under realistic incident pressure. "
            "Preview the hourglass flow: problem, design logic, system, demo, results, and close."
        ),
    },
    {
        "title": "The Hook: What Does a TPM Actually Do?",
        "notes": (
            "Use this diagram to ground the role reality: a TPM is not doing one task at a time. "
            "They coordinate incidents, architecture reviews, risk, and communication across teams. "
            "Set up the motivation: new TPMs are expected to reason across this complexity immediately."
        ),
    },
    {
        "title": "Motivation & Research Questions",
        "notes": (
            "Frame the problem as a learning systems gap, not an effort gap. "
            "Read RQ1 through RQ3 briefly and connect each to a measurable platform behavior: "
            "reasoning quality, unsupported claim reduction, and retrieval quality improvement."
        ),
    },
    {
        "title": "The Problem: Architecture Blindness",
        "notes": (
            "Walk through Nancy as a concrete learner persona. "
            "Highlight the three reinforcing gaps: overload, anxiety, and social isolation. "
            "Close with the key point that onboarding resources alone do not break this cycle."
        ),
    },
    {
        "title": "From Research to Design — Every Feature Traces Back",
        "notes": (
            "This is the research-to-implementation traceability slide. "
            "Move row by row and show that each major feature has a direct learning science rationale. "
            "Emphasize this is a deliberate design process, not feature accumulation."
        ),
    },
    {
        "title": "System Architecture",
        "notes": (
            "Read left to right: learner interaction, frontend, backend orchestration, then sub-systems. "
            "Point out where retrieval, graph context, and coaching logic connect. "
            "End with evidence gating and testing as reliability controls, not afterthoughts."
        ),
    },
    {
        "title": "The Four-Step Practice Loop",
        "notes": (
            "Explain the pedagogical sequencing: brief for context, investigate for hypothesis-building, "
            "decide for accountable reasoning, then feedback for correction. "
            "Call out that evidence gating is enforced before feedback to prevent answer hunting."
        ),
    },
    {
        "title": "Evidence Gating & AI Coaching",
        "notes": (
            "Differentiate two controls: evidence gating validates claims, and the coach shapes reflection. "
            "State clearly that the coach is Socratic and step-aware, and does not reveal final answers. "
            "Mention local model execution to reinforce privacy and offline operation."
        ),
    },
    {
        "title": "Graph-Augmented Retrieval",
        "notes": (
            "Describe pipeline flow quickly, then spend time on the bar chart. "
            "Use the ablation sequence to show incremental value: baseline, graph context, then evidence gating. "
            "Anchor the takeaway in measured improvement and unsupported-claim elimination."
        ),
    },
    {
        "title": "Synthetic Scenario Corpus — Omni-Mart",
        "notes": (
            "Show coverage and realism: four domains, three scenarios each, incident-style structure. "
            "Reinforce that this is fully synthetic to avoid sensitive data exposure. "
            "Connect corpus design to transferability for future role and domain expansion."
        ),
    },
    {
        "title": "Live Tool Demonstration",
        "notes": (
            "Use this as a transition slide and narrate the demo path before switching screens. "
            "Keep the demo focused on one scenario and the full four-step cycle. "
            "Explicitly show a gating event and one AI coaching interaction."
        ),
    },
    {
        "title": "Evaluation Results",
        "notes": (
            "Tie each result card back to RQ1 through RQ3 and keep claims precise. "
            "State the score lift from 0.856 to 0.963 and mention complete unsupported-claim elimination. "
            "Close with instrumentation and test coverage as confidence evidence."
        ),
    },
    {
        "title": "Key Contributions",
        "notes": (
            "Summarize contributions as an integrated learning system: structured practice, evidence controls, "
            "guardrailed coaching, and graph-aware retrieval. "
            "Avoid repetition of implementation details and emphasize learner impact."
        ),
    },
    {
        "title": "Future Work & Conclusion",
        "notes": (
            "Present future work as concrete next validation and scale steps: user study, larger corpus, "
            "production retrieval, and collaborative modes. "
            "Finish with the core thesis: architecture fluency can be taught through guided practice."
        ),
    },
    {
        "title": "References",
        "notes": (
            "Keep this brief while references are visible. "
            "Thank the audience, restate the educational technology focus, and invite questions on "
            "method, evaluation, or deployment roadmap."
        ),
    },
]


def _write_speaker_notes_doc(path):
    lines = [
        "# OmniMentor Final Presentation Speaker Notes",
        "",
        "This narration guide mirrors the 15-slide deck and is generated from scripts/gen_pptx_v2.py.",
        "",
    ]
    for idx, item in enumerate(SPEAKER_NOTES, start=1):
        lines.append(f"## Slide {idx}: {item['title']}")
        lines.append("")
        lines.append(item["notes"])
        lines.append("")

    with open(path, "w", encoding="utf-8") as f:
        f.write("\n".join(lines).rstrip() + "\n")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title / Hook opener
# ═══════════════════════════════════════════════════════════════════════════
def s01_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)

    _tb(slide, 1.0, 1.9, 8.0, 1.6, "OmniMentor",
        sz=70, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    div = slide.shapes.add_shape(1, Inches(3.5), Inches(3.55), Inches(3.0), Inches(0.05))
    div.fill.solid(); div.fill.fore_color.rgb = BLUE; div.line.fill.background()

    _tb(slide, 1.0, 3.7, 8.0, 1.1,
        "Guided scenario-based practice for developing\narchitecture reasoning in new Technical Program Managers.",
        sz=21, color=DIM, align=PP_ALIGN.CENTER)
    _tb(slide, 1.0, 4.9, 8.0, 0.55,
        "Arvind Sharma  ·  Georgia Institute of Technology  ·  CS 6460  ·  Spring 2026",
        sz=14, color=BLUE, align=PP_ALIGN.CENTER)
    _tb(slide, 1.0, 5.5, 8.0, 0.45,
        "Educational Technology  ·  Development Track",
        sz=13, color=DIM, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Hook: What Does a TPM Actually Do?
# ═══════════════════════════════════════════════════════════════════════════
def s02_hook_tpm(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _title(slide, "The Hook: What Does a TPM Actually Do?")

    # Central TPM oval
    _oval(slide, 4.3, 3.2, 1.4, 1.0, PURPLE, border=WHITE, label="TPM", lsz=16)

    # 6 surrounding role nodes with connecting lines
    nodes = [
        (1.6,  1.7,  "Incident\nResponse"),
        (4.3,  1.2,  "Architecture\nReview"),
        (7.0,  1.7,  "Cross-team\nCoordination"),
        (1.6,  4.9,  "Risk &\nMitigation"),
        (4.3,  5.4,  "Stakeholder\nUpdates"),
        (7.0,  4.9,  "Delivery\nTracking"),
    ]
    tpm_cx, tpm_cy = 5.0, 3.7   # center of TPM oval

    for nx, ny, label in nodes:
        ncx = nx + 0.7; ncy = ny + 0.42
        # thin connector
        lx = min(ncx, tpm_cx); ly = min(ncy, tpm_cy)
        rw = abs(ncx - tpm_cx); rh = abs(ncy - tpm_cy)
        if rw < 0.05: rw = 0.05
        if rh < 0.05: rh = 0.05
        conn = slide.shapes.add_shape(1, Inches(lx), Inches(ly), Inches(rw), Inches(rh))
        conn.fill.solid(); conn.fill.fore_color.rgb = RGBColor(0x1E, 0x30, 0x50)
        conn.line.color.rgb = RGBColor(0x3B, 0x82, 0xF6); conn.line.width = Pt(0.75)

        _oval(slide, nx, ny, 1.4, 0.85, CARD, border=BLUE, label=label, lsz=10, lcol=DIM)

    _tb(slide, 0.4, 6.5, 9.2, 0.5,
        "A TPM spans 30+ microservices — incident bridges, architecture reviews, delivery tracking — simultaneously, on day one.",
        sz=13, color=DIM, align=PP_ALIGN.CENTER)
    _footer(slide)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Motivation & Research Questions
# ═══════════════════════════════════════════════════════════════════════════
def s03_motivation(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _title(slide, "Motivation & Research Questions")

    _tb(slide, 0.4, 1.05, 9.2, 0.45,
        "New TPMs have no structured way to build architecture reasoning. Documentation doesn't teach. Chat gives answers, not understanding.",
        sz=14, color=DIM)

    rqs = [
        ("RQ1", "Does scenario-based practice improve architecture reasoning scores?", BLUE),
        ("RQ2", "Does claim-level evidence gating reduce unsupported reasoning?", PURPLE),
        ("RQ3", "Does graph-augmented retrieval outperform flat RAG on reasoning quality?", GREEN),
    ]
    for i, (rq, text, col) in enumerate(rqs):
        t = 1.75 + i * 1.55
        _rect(slide, 0.4, t, 9.2, 1.3, CARD, border=col)
        _tb(slide, 0.65, t + 0.1, 1.0, 1.1, rq, sz=22, color=col, bold=True, align=PP_ALIGN.CENTER)
        _tb(slide, 1.8, t + 0.25, 7.5, 0.85, text, sz=16, color=WHITE)

    _footer(slide)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — The Problem: Three Reinforcing Gaps
# ═══════════════════════════════════════════════════════════════════════════
def s04_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _title(slide, "The Problem: Architecture Blindness")

    _tb(slide, 0.4, 1.0, 9.2, 0.48,
        "Nancy — new TPM, Omni-Mart — 30+ microservices, day-one incident meeting, zero mental model.",
        sz=14, color=DIM)

    cards = [
        ("🧠", "Cognitive Overload",
         "Too many systems, too many names —\nworking memory maxed out.",
         "Sweller, 1988", BLUE),
        ("😨", "Performance Anxiety",
         "Fear of asking 'obvious' questions leads\nto hesitation and avoidance.",
         "Bandura, 1977", PURPLE),
        ("👥", "Social Isolation",
         "In the meeting but can't contribute —\nstuck at the periphery.",
         "Lave & Wenger, 1991", GREEN),
    ]
    for i, (icon, heading, body, cite, col) in enumerate(cards):
        l = 0.35 + i * 3.15
        _rect(slide, l, 1.65, 2.95, 4.3, CARD, border=col)
        _tb(slide, l + 0.1, 1.78, 2.75, 0.6, icon,   sz=32, align=PP_ALIGN.CENTER)
        _tb(slide, l + 0.1, 2.5,  2.75, 0.55, heading, sz=16, color=col, bold=True, align=PP_ALIGN.CENTER)
        _tb(slide, l + 0.1, 3.15, 2.75, 1.5, body,   sz=13, color=WHITE, align=PP_ALIGN.CENTER)
        _tb(slide, l + 0.1, 5.05, 2.75, 0.4, cite,   sz=11, color=DIM, italic=True, align=PP_ALIGN.CENTER)

    _callout(slide, 0.4, 6.15, 9.2, 0.62,
             "Without structured practice, these three gaps reinforce each other — persisting well beyond onboarding.")
    _footer(slide)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Literature → Design Mapping
# ═══════════════════════════════════════════════════════════════════════════
def s05_lit(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _title(slide, "From Research to Design — Every Feature Traces Back")

    rows = [
        ("Cognitive Load Theory",          "Sweller, 1988; Paas & van Merriënboer, 1994",
         "Four-step practice loop: Brief → Investigate → Decide → Feedback"),
        ("Situated Cognition & LPP",       "Brown et al., 1989; Lave & Wenger, 1991",
         "12 authentic incident scenarios across 4 Omni-Mart domains"),
        ("Intelligent Tutoring",           "VanLehn, 2011; Chi et al., 1989",
         "Claim-level evidence gating — unsupported claims flagged before feedback"),
        ("Design-Based Research",          "Barab & Squire, 2004; DBR Collective, 2003",
         "Weekly design-evaluate-refine cycles with ablation validation"),
    ]
    hdrs   = ["Theory",  "Key Sources",  "Design Decision"]
    col_l  = [0.4,  2.75,  6.05]
    col_w  = [2.25, 3.2,   3.55]

    for j, (hdr, lft, wd) in enumerate(zip(hdrs, col_l, col_w)):
        _rect(slide, lft, 1.1, wd, 0.42, RGBColor(0x1A, 0x26, 0x40),
              border=RGBColor(0x2A, 0x3A, 0x58),
              label=hdr, lsz=12, lcol=BLUE, lbold=True)

    for i, (theory, sources, decision) in enumerate(rows):
        t    = 1.62 + i * 1.25
        fill = CARD if i % 2 == 0 else CARD2
        data = [theory, sources, decision]
        for j, (val, lft, wd) in enumerate(zip(data, col_l, col_w)):
            _rect(slide, lft, t, wd, 1.18, fill,
                  border=RGBColor(0x1E, 0x2E, 0x48))
            _tb(slide, lft + 0.1, t + 0.1, wd - 0.2, 1.0,
                val, sz=12, color=WHITE if j < 2 else DIM)

    _footer(slide)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — System Architecture (block diagram)
# ═══════════════════════════════════════════════════════════════════════════
def s06_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _title(slide, "System Architecture")

    # ── Tier 1: Learner → React → FastAPI ───────────────────────────────
    _rect(slide, 0.3,  2.0, 1.6, 0.85, CARD,
          border=BLUE, label="Learner\n(Browser)", lsz=11, lcol=DIM)
    _arrow_r(slide, 1.92, 2.28)
    _rect(slide, 2.28, 1.72, 1.9, 1.4, RGBColor(0x0D, 0x1A, 0x33),
          border=BLUE, label="React\nFrontend\n(Vite + TS)", lsz=11, lcol=BLUE)
    _arrow_r(slide, 4.2, 2.28)
    _rect(slide, 4.55, 1.72, 1.9, 1.4, RGBColor(0x12, 0x10, 0x30),
          border=PURPLE, label="FastAPI\nBackend\n(Python)", lsz=11, lcol=PURPLE)

    # ── Down arrow ───────────────────────────────────────────────────────
    _arrow_d(slide, 5.3, 3.14)

    # ── Tier 2: three sub-systems ─────────────────────────────────────────
    subs = [
        (3.25, "Scenario\nCorpus\n(JSON)",         BLUE),
        (5.05, "Graph-RAG\n(NetworkX)",             PURPLE),
        (6.85, "Ollama LLM\nCoach\n(Llama 3.2)",   GREEN),
    ]
    for lft, label, col in subs:
        _rect(slide, lft, 3.55, 1.65, 1.35, CARD, border=col,
              label=label, lsz=10, lcol=col)

    # thin vertical drops from FastAPI to sub boxes
    for cx in [4.08, 5.0, 5.93]:
        ln = slide.shapes.add_shape(1, Inches(cx), Inches(3.13), Inches(0.02), Inches(0.42))
        ln.fill.solid(); ln.fill.fore_color.rgb = BLUE; ln.line.fill.background()

    # ── Tier 3: evaluation + testing ─────────────────────────────────────
    _rect(slide, 3.25, 5.15, 5.25, 0.78,
          RGBColor(0x0A, 0x12, 0x28),
          border=RGBColor(0x2E, 0x44, 0x60),
          label="Evidence Gating  ·  Claim Scorer  ·  Playwright E2E (61 tests)",
          lsz=11, lcol=DIM)

    _tb(slide, 0.4, 6.18, 9.2, 0.48,
        "Tech stack:  TypeScript  ·  React  ·  FastAPI  ·  Python  ·  NetworkX  ·  Ollama  ·  Playwright",
        sz=12, color=DIM, align=PP_ALIGN.CENTER)
    _footer(slide)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Four-Step Practice Loop (flowchart)
# ═══════════════════════════════════════════════════════════════════════════
def s07_flowchart(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _title(slide, "The Four-Step Practice Loop")

    steps = [
        ("1  Brief",       "Read the incident scenario and system context snapshot.",           BLUE),
        ("2  Investigate",  "Identify affected services; build a hypothesis about the failure.", PURPLE),
        ("3  Decide",       "Propose root cause + resolution. Every claim requires evidence.",   GREEN),
        ("4  Feedback",     "AI coach reveals reasoning gaps — Socratic, never answer-giving.",  AMBER),
    ]

    for i, (label, desc, col) in enumerate(steps):
        t = 1.25 + i * 1.42
        # step box
        _rect(slide, 0.4, t, 1.7, 1.0, CARD, border=col,
              label=label, lsz=13, lcol=col, lbold=True)
        # description
        _tb(slide, 2.25, t + 0.1, 7.3, 0.82, desc, sz=15, color=WHITE)
        # down arrow between steps
        if i < 3:
            _arrow_d(slide, 1.15, t + 1.02, w=0.48, h=0.35)

    # Loop annotation on right
    _tb(slide, 8.55, 1.25, 1.1, 5.5,
        "Design\nEvaluate\nRefine\n\n(DBR\nLoop)",
        sz=11, color=DIM, align=PP_ALIGN.CENTER)
    # right-side brace line
    ln = slide.shapes.add_shape(1, Inches(8.48), Inches(1.25), Inches(0.04), Inches(5.0))
    ln.fill.solid(); ln.fill.fore_color.rgb = RGBColor(0x2A, 0x40, 0x60); ln.line.fill.background()

    _callout(slide, 0.4, 7.02, 9.2, 0.46,
             "Evidence gating fires at Step 3 — unsupported claims blocked before the learner reaches feedback.")
    _footer(slide)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Evidence Gating & AI Coach
# ═══════════════════════════════════════════════════════════════════════════
def s08_gating_coach(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _title(slide, "Evidence Gating & AI Coaching")

    # Left card
    _rect(slide, 0.4, 1.1, 4.4, 5.7, CARD, border=PURPLE)
    _tb(slide, 0.6, 1.2, 4.0, 0.52, "Evidence Gating", sz=18, color=PURPLE, bold=True)
    gating = [
        "Learner submits a claim during Step 3",
        "System checks: is this backed by retrieved evidence?",
        "Unsupported → hint shown, not the answer",
        "Supported → proceeds to AI coaching step",
        "Prevents answer-fishing through the feedback loop",
    ]
    for j, b in enumerate(gating):
        _tb(slide, 0.6, 1.88 + j * 0.92, 4.1, 0.8, f"• {b}", sz=13, color=WHITE)

    # Right card
    _rect(slide, 5.2, 1.1, 4.4, 5.7, CARD, border=BLUE)
    _tb(slide, 5.4, 1.2, 4.0, 0.52, "AI Coach Guardrails", sz=18, color=BLUE, bold=True)
    coach = [
        "Step-aware system prompt — knows learner's current step",
        "Never reveals the answer directly",
        "Asks Socratic questions to surface reasoning gaps",
        "Off-topic or harmful queries: politely redirected",
        "Local LLM (Ollama / Llama 3.2) — no data leaves the machine",
    ]
    for j, b in enumerate(coach):
        _tb(slide, 5.4, 1.88 + j * 0.92, 4.1, 0.8, f"• {b}", sz=13, color=WHITE)

    _footer(slide)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Graph-Augmented Retrieval + Ablation Bar Chart
# ═══════════════════════════════════════════════════════════════════════════
def s09_retrieval(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _title(slide, "Graph-Augmented Retrieval")

    # Pipeline boxes
    pipe = [
        (0.35,  "Learner\nQuery",                BLUE),
        (2.25,  "BM25\nKeyword\n(Flat RAG)",     DIM),
        (4.15,  "Service\nDependency\nGraph",     PURPLE),
        (6.05,  "Re-ranked\nContext",             GREEN),
        (7.95,  "Ollama\nCoach",                 BLUE),
    ]
    for lft, label, col in pipe:
        _rect(slide, lft, 2.0, 1.72, 1.3, CARD, border=col,
              label=label, lsz=10, lcol=col)
        if lft < 7.95:
            _arrow_r(slide, lft + 1.74, 2.5)

    # Bar chart header
    _tb(slide, 0.4, 3.55, 9.2, 0.4,
        "Ablation study — Mean reasoning score  (higher is better)",
        sz=13, color=DIM, bold=True, align=PP_ALIGN.LEFT)

    # Three bars: baseline, +graph, +gating
    # Bar area: x starts at 2.7, max width = 6.5
    bar_data = [
        ("Flat RAG (baseline)",       0.856, DIM),
        ("+ Graph context",           0.921, BLUE),
        ("+ Evidence gating (final)", 0.963, GREEN),
    ]
    max_w = 6.3
    for k, (label, score, col) in enumerate(bar_data):
        t = 4.15 + k * 0.95
        filled_w = (score / 1.0) * max_w
        _hbar(slide, 2.65, t, max_w, filled_w, 0.62, col)
        _tb(slide, 0.4, t + 0.05, 2.2, 0.55, label, sz=11, color=DIM)
        _tb(slide, 2.65 + filled_w + 0.08, t + 0.05, 0.9, 0.55,
            f"{score:.3f}", sz=13, color=col, bold=True)

    _callout(slide, 0.4, 7.05, 9.2, 0.48,
             "+12.5% improvement over flat RAG baseline  ·  All unsupported claims eliminated",
             border=GREEN)
    _footer(slide)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Synthetic Scenario Corpus
# ═══════════════════════════════════════════════════════════════════════════
def s10_corpus(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _title(slide, "Synthetic Scenario Corpus — Omni-Mart")

    _tb(slide, 0.4, 1.05, 9.2, 0.45,
        "12 realistic incident scenarios across 4 domains. Fully synthetic — no real PII, no proprietary data.",
        sz=14, color=DIM)

    domains = [
        ("🛒", "Cart &\nCheckout",       "3 scenarios\nPayment service failure,\ncart state corruption",   BLUE),
        ("📦", "Fulfillment &\nLogistics","3 scenarios\nWarehouse sync issues,\ndelivery pipeline stalls",  PURPLE),
        ("🗂️", "Product\nCatalog",       "3 scenarios\nSearch index staleness,\ncatalog write conflicts",   GREEN),
        ("🔒", "Risk &\nCompliance",      "3 scenarios\nFraud signal delays,\nrate-limiting failures",      AMBER),
    ]
    for i, (icon, name, detail, col) in enumerate(domains):
        l = 0.35 + i * 2.35
        _rect(slide, l, 1.7, 2.2, 4.5, CARD, border=col)
        _tb(slide, l + 0.1, 1.82, 2.0, 0.62, icon,   sz=28, align=PP_ALIGN.CENTER)
        _tb(slide, l + 0.1, 2.55, 2.0, 0.65, name,   sz=14, color=col, bold=True, align=PP_ALIGN.CENTER)
        _tb(slide, l + 0.1, 3.35, 2.0, 2.5,  detail, sz=11, color=DIM, align=PP_ALIGN.CENTER)

    _tb(slide, 0.4, 6.45, 9.2, 0.48,
        "Each scenario: title · symptom · 5-8 service nodes · evidence chunks · expected root cause",
        sz=13, color=DIM, align=PP_ALIGN.CENTER)
    _footer(slide)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Tool Demo Transition
# ═══════════════════════════════════════════════════════════════════════════
def s11_demo(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)

    _tb(slide, 1.0, 1.7, 8.0, 1.2, "Live Tool Demonstration",
        sz=46, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    _tb(slide, 1.0, 2.9, 8.0, 0.58,
        "[ Switching to screen recording ]",
        sz=20, color=BLUE, align=PP_ALIGN.CENTER)

    waypoints = [
        "Scenario selection  →  pick an incident from the corpus",
        "Four-step flow  →  Brief, Investigate, Decide, Feedback",
        "Evidence gating  →  submit a claim and watch it fire",
        "AI coaching assistant  →  Socratic prompts in action",
        "System dependency graph  →  service topology explorer",
    ]
    for j, item in enumerate(waypoints):
        _tb(slide, 2.0, 3.72 + j * 0.6, 6.5, 0.52,
            f"▸  {item}", sz=14, color=DIM)

    _footer(slide)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Evaluation Results
# ═══════════════════════════════════════════════════════════════════════════
def s12_results(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _title(slide, "Evaluation Results")

    rqs = [
        ("RQ1", "Architecture reasoning scores improved across all four practice steps.",           BLUE),
        ("RQ2", "Evidence gating eliminated all unsupported claims in final submissions.",          GREEN),
        ("RQ3", "Graph-augmented retrieval lifted scores 0.856 → 0.963  (+12.5%)",                 PURPLE),
    ]
    for i, (rq, finding, col) in enumerate(rqs):
        t = 1.2 + i * 1.58
        _rect(slide, 0.4, t, 9.2, 1.38, CARD, border=col)
        _tb(slide, 0.65, t + 0.12, 1.0, 1.1, rq, sz=22, color=col, bold=True, align=PP_ALIGN.CENTER)
        _tb(slide, 1.82, t + 0.25, 6.7, 0.88, finding, sz=15, color=WHITE)
        _tb(slide, 8.65, t + 0.35, 0.75, 0.5, "✓", sz=22, color=col, bold=True, align=PP_ALIGN.CENTER)

    _callout(slide, 0.4, 6.05, 9.2, 0.62,
             "61 Playwright E2E tests  ·  All passing  ·  Full RQ1–RQ3 instrumentation  ·  Ablation validated",
             border=GREEN)
    _footer(slide)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Key Contributions
# ═══════════════════════════════════════════════════════════════════════════
def s13_contributions(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _title(slide, "Key Contributions")

    contribs = [
        ("Practice\nFramework",  "Four-step scenario loop decomposes architecture reasoning into teachable, measurable sub-skills.",  BLUE),
        ("Evidence\nGating",     "Claim-level validation surfaces unsupported reasoning gaps before the learner sees any answer.",     PURPLE),
        ("AI Coach",             "Local LLM (Ollama / Llama 3.2) — step-aware guardrails, Socratic-only, runs fully offline.",        GREEN),
        ("Graph-RAG\nPipeline",  "Service dependency graph re-ranks retrieval context, lifting scores +12.5% over flat baseline.",     AMBER),
    ]
    for i, (title, body, col) in enumerate(contribs):
        t = 1.2 + i * 1.4
        _rect(slide, 0.4, t, 9.2, 1.25, CARD, border=col)
        _tb(slide, 0.65, t + 0.08, 1.7, 1.1, title, sz=15, color=col, bold=True, align=PP_ALIGN.CENTER)
        _tb(slide, 2.5,  t + 0.2,  6.8, 0.88, body,  sz=14, color=WHITE)

    _footer(slide)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Future Work & Conclusion
# ═══════════════════════════════════════════════════════════════════════════
def s14_future(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _title(slide, "Future Work & Conclusion")

    # Left
    _rect(slide, 0.4, 1.15, 4.5, 5.65, CARD, border=BLUE)
    _tb(slide, 0.6, 1.25, 4.1, 0.52, "Future Work", sz=18, color=BLUE, bold=True)
    fw = [
        "Controlled user study with real TPMs",
        "Expanded corpus — additional domains",
        "Production retrieval (Qdrant + Neo4j)",
        "Multi-role support (SREs, engineers)",
        "Collaborative practice sessions",
    ]
    for j, f in enumerate(fw):
        _tb(slide, 0.6, 1.9 + j * 0.9, 4.2, 0.82, f"• {f}", sz=13, color=WHITE)

    # Right
    _rect(slide, 5.2, 1.15, 4.5, 5.65, CARD, border=PURPLE)
    _tb(slide, 5.4, 1.25, 4.1, 0.52, "Conclusion", sz=18, color=PURPLE, bold=True)
    _tb(slide, 5.4, 1.9, 4.1, 4.6,
        "Architecture fluency is a learnable skill,\nnot a documentation-access problem.\n\n"
        "OmniMentor gives new TPMs a low-stakes practice space — "
        "guided by research, instrumented for evaluation, "
        "and scaffolded by a guardrailed AI coach that nudges without giving answers.",
        sz=13, color=WHITE, wrap=True)

    _footer(slide)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 15 — References (two columns)
# ═══════════════════════════════════════════════════════════════════════════
def s15_refs(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _title(slide, "References", sz=28)

    col1 = [
        "1. Barab, S. & Squire, K. (2004). Design-Based Research:\n   Putting a stake in the ground. J. Learning Sciences, 13(1), 1-14.",
        "2. Bandura, A. (1977). Self-efficacy: Toward a unifying\n   theory of behavioral change. Psychological Review, 84(2), 191-215.",
        "3. Brown, J. S., Collins, A., & Duguid, P. (1989). Situated\n   cognition and the culture of learning. Ed. Researcher, 18(1), 32-42.",
        "4. Chi, M. T. H. et al. (1989). Self-explanations: How\n   students study examples. Cognitive Science, 13(2), 145-182.",
        "5. Collins, A., Brown, J. S., & Newman, S. E. (1989).\n   Cognitive apprenticeship. Lawrence Erlbaum.",
        "6. DBR Collective. (2003). Design-Based Research:\n   An emerging paradigm. Educational Researcher, 32(1), 5-8.",
    ]
    col2 = [
        "7.  Lave, J., & Wenger, E. (1991). Situated Learning:\n    Legitimate Peripheral Participation. Cambridge Univ. Press.",
        "8.  Paas, F. & van Merriënboer, J. J. G. (1994). Instructional\n    control of cognitive load. Ed. Psychology Review, 6(4), 351-371.",
        "9.  Sweller, J. (1988). Cognitive load during problem solving.\n    Cognitive Science, 12(2), 257-285.",
        "10. VanLehn, K. (2011). The relative effectiveness of human\n    tutoring and ITS. Educational Psychologist, 46(4), 197-221.",
        "11. Beyer, B. et al. (2016). Site Reliability Engineering.\n    O'Reilly Media.",
        "12. Marteau, T. M. & Bekker, H. (1992). STAI short-form\n    development. Brit. J. Clinical Psychology, 31(3), 301-306.",
    ]
    for j, r in enumerate(col1):
        _tb(slide, 0.4, 1.1 + j * 1.0, 4.6, 0.92, r, sz=10, color=DIM)
    for j, r in enumerate(col2):
        _tb(slide, 5.1, 1.1 + j * 1.0, 4.6, 0.92, r, sz=10, color=DIM)

    _footer(slide)


# ═══════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════
def main():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.6)

    s01_title(prs)
    s02_hook_tpm(prs)
    s03_motivation(prs)
    s04_problem(prs)
    s05_lit(prs)
    s06_architecture(prs)
    s07_flowchart(prs)
    s08_gating_coach(prs)
    s09_retrieval(prs)
    s10_corpus(prs)
    s11_demo(prs)
    s12_results(prs)
    s13_contributions(prs)
    s14_future(prs)
    s15_refs(prs)

    for slide, note in zip(prs.slides, SPEAKER_NOTES):
        _set_notes(slide, note["notes"])

    out_dir = "personal/prep-final-video"
    os.makedirs(out_dir, exist_ok=True)

    out_pptx = f"{out_dir}/OMNIMENTOR_FINAL_PRESENTATION_FROM_INSTRUCTIONS.pptx"
    out_notes = f"{out_dir}/OMNIMENTOR_FINAL_PRESENTATION_SPEAKER_NOTES.md"

    prs.save(out_pptx)
    _write_speaker_notes_doc(out_notes)
    print(f"Done — {len(prs.slides)} slides saved to {out_pptx}")
    print(f"Done — speaker notes doc saved to {out_notes}")

if __name__ == "__main__":
    main()
