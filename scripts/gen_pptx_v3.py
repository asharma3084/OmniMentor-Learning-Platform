"""
OmniMentor — CS 6460 Final Presentation Generator v3
8 slides, aligned with official GT CS6460 presentation guidance.

Time budget:  Slides 1-5 (~2 min) → Live demo (~7 min) → Slides 6-8 (~1 min)
Design:       Dark navy, minimal text, presenter narrates, demo is the star.
Speaker notes: Goal 3 — seed conversation, invite feedback. Talking points, not script.
"""
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ─────────────────────────────────────────────────────────────────
BG      = RGBColor(0x06, 0x0B, 0x18)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BLUE    = RGBColor(0x3B, 0x82, 0xF6)
PURPLE  = RGBColor(0x8B, 0x5C, 0xF6)
GREEN   = RGBColor(0x10, 0xB9, 0x81)
AMBER   = RGBColor(0xF5, 0x9E, 0x0B)
DIM     = RGBColor(0xA0, 0xAE, 0xC0)
CARD    = RGBColor(0x11, 0x18, 0x27)
CARD2   = RGBColor(0x0D, 0x14, 0x22)
DARK2   = RGBColor(0x0A, 0x12, 0x24)


# ── Low-level helpers ────────────────────────────────────────────────────────
def _bg(slide):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = BG


def _tb(slide, l, t, w, h, text, sz=15, color=WHITE, bold=False,
        italic=False, align=PP_ALIGN.LEFT, wrap=True):
    s = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    s.text_frame.word_wrap = wrap
    p = s.text_frame.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    return s


def _rect(slide, l, t, w, h, fill, border=None, label="",
          lsz=14, lcol=WHITE, lbold=False, wrap=True, desc=""):
    s = slide.shapes.add_shape(
        1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(1.2)
    else:
        s.line.fill.background()
    if label:
        tf = s.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.text = label
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(lsz)
        p.font.color.rgb = lcol
        p.font.bold = lbold
    if desc:
        s.name = desc
    return s


def _oval(slide, l, t, w, h, fill, border=WHITE, label="", lsz=14, lcol=WHITE, desc=""):
    s = slide.shapes.add_shape(
        9, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = border
    s.line.width = Pt(1)
    if label:
        tf = s.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(lsz)
        p.font.color.rgb = lcol
        p.font.bold = True
    if desc:
        s.name = desc
    return s


def _arrow_r(slide, l, t, w=0.32, h=0.22):
    s = slide.shapes.add_shape(13, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = BLUE; s.line.fill.background()
    s.name = "Arrow right"


def _arrow_d(slide, l, t, w=0.22, h=0.32):
    s = slide.shapes.add_shape(36, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = BLUE; s.line.fill.background()
    s.name = "Arrow down"


def _title(slide, text, sz=33):
    _tb(slide, 0.4, 0.18, 9.2, 0.72, text, sz=sz, color=WHITE, bold=True)
    bar = slide.shapes.add_shape(
        1, Inches(0.4), Inches(0.88), Inches(9.2), Inches(0.04))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()


def _footer(slide):
    _tb(slide, 0.4, 7.18, 9.2, 0.32,
        "Arvind Sharma  ·  CS 6460  ·  Georgia Institute of Technology  ·  Spring 2026",
        sz=11, color=BLUE, align=PP_ALIGN.CENTER)


def _callout(slide, l, t, w, h, text, sz=14, border=BLUE):
    s = _rect(slide, l, t, w, h, DARK2, border=border, desc="Callout")
    tf = s.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(sz)
    p.font.color.rgb = WHITE


def _hbar(slide, l, t, total_w, filled_w, h, fill, label=""):
    bg = _rect(slide, l, t, total_w, h, RGBColor(0x1A, 0x24, 0x3A))
    fg = _rect(slide, l, t, filled_w, h, fill)
    if label:
        bg.name = label
        fg.name = f"{label} filled"


def _set_notes(slide, text):
    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    notes.text = text


# ── Speaker notes (Goal 3: seed conversation, invite feedback) ──────────
SPEAKER_NOTES = [
    {
        "title": "OmniMentor",
        "notes": (
            "OmniMentor is a guided practice environment that helps new Technical Program Managers "
            "build architecture reasoning through realistic incident scenarios.\n\n"
            "Quick roadmap: I'll set up the problem and design logic in about two minutes, "
            "then switch to a live demo for seven minutes, and close with results and next steps."
        ),
    },
    {
        "title": "The Problem: Architecture Blindness",
        "notes": (
            "Meet Nancy — new TPM, first day at an enterprise with thousands of microservices. "
            "She has access to everything: wikis, Jira boards, Confluence pages, architecture diagrams, "
            "even Coursera and LinkedIn Learning courses. But right now she's on her first incident bridge "
            "and none of that helps — because reference material doesn't build reasoning skill.\n\n"
            "This is what learning science calls cognitive overload — she must process system topology, "
            "service ownership, dependency chains, and governance constraints all at once, far exceeding "
            "working memory capacity (Sweller, 1988). She's afraid to ask a question that sounds obvious "
            "to everyone else (Bandura, 1977). And she's stuck at the periphery — in the meeting but unable "
            "to contribute until she builds a working mental model (Lave & Wenger, 1991).\n\n"
            "These three gaps reinforce each other and extend the ramp-up far beyond "
            "what any amount of documentation, courses, or wikis can address. "
            "OmniMentor targets that window — giving new TPMs a way to practice reasoning "
            "before the pressure is real."
        ),
    },
    {
        "title": "From Research to Design",
        "notes": (
            "Every major feature traces directly to a learning science theory — this wasn't built "
            "by feature accumulation.\n\n"
            "Cognitive Load Theory gives us the four-step decomposition. Situated Cognition gives us "
            "realistic incident scenarios. Intelligent Tutoring research gives us evidence gating. "
            "Design-Based Research gives us the iterative validation process.\n\n"
            "Plus: a local LLM coach (Ollama / Llama 3.2) that guides without revealing answers."
        ),
    },
    {
        "title": "System Architecture",
        "notes": (
            "Read left to right: learner in browser, React frontend, Express API backend orchestrating "
            "three sub-systems — the scenario corpus, graph-augmented retrieval via NetworkX, "
            "and the Ollama coaching LLM.\n\n"
            "Evidence gating and Playwright E2E tests sit underneath as reliability controls. "
            "The whole stack runs locally — no data leaves the machine."
        ),
    },
    {
        "title": "The Gap — Demo Transition",
        "notes": (
            "Think about what organizations already provide: wikis with hundreds of pages, Jira boards, "
            "Confluence spaces, architecture diagrams, even platforms like Coursera and LinkedIn Learning. "
            "New TPMs are drowning in reference material — but none of it teaches them to reason through "
            "a live incident, to figure out which service is upstream, to assess blast radius under pressure.\n\n"
            "That's the gap: the difference between having information available and being able to "
            "use it when it matters. Decades of learning science research — Sweller, Lave and Wenger, "
            "the cognitive apprenticeship tradition — all point to the same conclusion: "
            "you build reasoning skill through structured practice, not through reading.\n\n"
            "OmniMentor fills that gap. Let me show you how.\n\n"
            "[Switch to browser — demo path: scenario selection, four-step flow, evidence gating "
            "event, AI coaching interaction, service graph explorer]"
        ),
    },
    {
        "title": "Results & Contributions",
        "notes": (
            "Key result: graph-augmented retrieval with evidence gating improved reasoning scores "
            "from 0.856 to 0.963 — a 12.5 percent improvement — and eliminated all unsupported claims. "
            "Validated across 12 scenarios and 4 domains.\n\n"
            "Five contributions: scenario-based practice framework, claim-level evidence gating, "
            "guardrailed AI coaching with a local LLM and eight behavioral guardrails, "
            "graph-augmented retrieval validated by ablation study, and full RQ1–RQ3 "
            "instrumentation with 61 end-to-end tests.\n\n"
            "Scoped limitation: ablation-validated with synthetic data, no human subjects study yet "
            "— that's the most important next step."
        ),
    },
    {
        "title": "Future Work & Conclusion",
        "notes": (
            "Next steps: controlled user study with real TPMs, expanded scenario corpus, production "
            "retrieval backends (Qdrant + Neo4j), multi-role support, collaborative practice.\n\n"
            "The core insight, backed by decades of learning science research: architecture fluency "
            "is a skill you build through practice, not a documentation problem you solve with more wikis. "
            "Organizations already invest heavily in knowledge bases and training platforms — "
            "what's been missing is a practice environment where new TPMs can fail safely "
            "and build reasoning skill before the real incident bridge.\n\n"
            "I'd especially welcome feedback on the user study design and how to structure it "
            "for ecological validity — that's the most important next step."
        ),
    },
    {
        "title": "References",
        "notes": (
            "Thank you. These are the 12 references grounding the design and evaluation.\n\n"
            "Happy to discuss method choices, evaluation approach, or deployment roadmap."
        ),
    },
]


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title (5 sec)
# ═══════════════════════════════════════════════════════════════════════════
def s01_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)

    _tb(slide, 1.0, 1.9, 8.0, 1.6, "OmniMentor",
        sz=70, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    div = slide.shapes.add_shape(
        1, Inches(3.5), Inches(3.55), Inches(3.0), Inches(0.05))
    div.fill.solid(); div.fill.fore_color.rgb = BLUE; div.line.fill.background()
    div.name = "Title divider"

    _tb(slide, 1.0, 3.7, 8.0, 1.1,
        "Guided scenario-based practice for developing\narchitecture reasoning "
        "in new Technical Program Managers.",
        sz=21, color=DIM, align=PP_ALIGN.CENTER)
    _tb(slide, 1.0, 4.9, 8.0, 0.55,
        "Arvind Sharma  ·  Georgia Institute of Technology  ·  CS 6460  ·  Spring 2026",
        sz=14, color=BLUE, align=PP_ALIGN.CENTER)
    _tb(slide, 1.0, 5.5, 8.0, 0.45,
        "Educational Technology  ·  Development Track",
        sz=14, color=DIM, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem: Architecture Blindness (40 sec)
# ═══════════════════════════════════════════════════════════════════════════
def s02_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _title(slide, "The Problem: Architecture Blindness")

    _tb(slide, 0.4, 1.05, 9.2, 0.5,
        "Nancy — new TPM at Omni-Mart. Thousands of microservices. She has the wikis, "
        "the Confluence pages, the training courses — but none of that builds reasoning skill.",
        sz=16, color=DIM)

    cards = [
        ("🧠", "Cognitive\nOverload",
         "Hundreds of services, ownership\nregistries, dependency chains —\nall competing for working memory.",
         "Sweller, 1988", BLUE),
        ("😨", "Performance\nAnxiety",
         "Fear of asking a question that\nsounds obvious to the senior\nengineers on the bridge.",
         "Bandura, 1977", PURPLE),
        ("👥", "Social\nIsolation",
         "Present in the meeting but unable\nto contribute — stuck at the\nperiphery.",
         "Lave & Wenger, 1991", GREEN),
    ]
    for i, (icon, heading, body, cite, col) in enumerate(cards):
        l = 0.35 + i * 3.15
        card_desc = f"Card: {heading.replace(chr(10), ' ')}"
        _rect(slide, l, 1.75, 2.95, 4.2, CARD, border=col, desc=card_desc)
        _tb(slide, l + 0.1, 1.9, 2.75, 0.65, icon,
            sz=36, align=PP_ALIGN.CENTER)
        _tb(slide, l + 0.1, 2.7, 2.75, 0.7, heading,
            sz=18, color=col, bold=True, align=PP_ALIGN.CENTER)
        _tb(slide, l + 0.1, 3.5, 2.75, 1.4, body,
            sz=15, color=WHITE, align=PP_ALIGN.CENTER)
        _tb(slide, l + 0.1, 5.1, 2.75, 0.4, cite,
            sz=12, color=DIM, italic=True, align=PP_ALIGN.CENTER)

    _callout(slide, 0.4, 6.2, 9.2, 0.6,
             "These gaps reinforce each other. Without structured practice, "
             "they persist well beyond onboarding.")
    _footer(slide)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Research → Design Mapping (40 sec)
# ═══════════════════════════════════════════════════════════════════════════
def s03_research_design(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _title(slide, "From Research to Design — Every Feature Traces Back")

    cols = [
        ("Cognitive\nLoad Theory",
         "Sweller, 1988\nPaas & van Merriënboer, 1994",
         "Four-step practice loop:\nBrief → Investigate →\nDecide → Feedback",
         BLUE),
        ("Situated\nCognition & LPP",
         "Brown et al., 1989\nLave & Wenger, 1991",
         "12 realistic incident\nscenarios across 4\nOmni-Mart domains",
         PURPLE),
        ("Intelligent\nTutoring",
         "VanLehn, 2011\nChi et al., 1989",
         "Claim-level evidence\ngating — unsupported\nclaims flagged",
         GREEN),
        ("Design-Based\nResearch",
         "Barab & Squire, 2004\nDBR Collective, 2003",
         "Weekly design-evaluate-\nrefine iterations with\nablation validation",
         AMBER),
    ]

    for i, (theory, cite, decision, col) in enumerate(cols):
        l = 0.25 + i * 2.45
        cw = 2.3

        # Theory box (top)
        _rect(slide, l, 1.1, cw, 1.55, CARD, border=col,
              desc=f"Theory: {theory.replace(chr(10), ' ')}")
        _tb(slide, l + 0.1, 1.18, cw - 0.2, 0.7, theory,
            sz=14, color=col, bold=True, align=PP_ALIGN.CENTER)
        _tb(slide, l + 0.1, 1.95, cw - 0.2, 0.6, cite,
            sz=11, color=DIM, italic=True, align=PP_ALIGN.CENTER)

        # Down arrow
        _arrow_d(slide, l + cw / 2 - 0.11, 2.7, w=0.22, h=0.35)

        # Design decision box (bottom)
        _rect(slide, l, 3.15, cw, 1.8, CARD, border=col,
              desc=f"Design: {decision.replace(chr(10), ' ')}")
        _tb(slide, l + 0.1, 3.25, cw - 0.2, 1.6, decision,
            sz=14, color=WHITE, align=PP_ALIGN.CENTER)

    _callout(slide, 0.4, 5.25, 9.2, 0.65,
             "+ AI coaching assistant: local LLM (Ollama / Llama 3.2) provides "
             "step-aware guidance without revealing answers.",
             border=BLUE)
    _footer(slide)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — System Architecture (20 sec)
# ═══════════════════════════════════════════════════════════════════════════
def s04_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _title(slide, "System Architecture")

    # ── Tier 1: Learner → React → Express API  (centered) ───────────────
    _rect(slide, 0.8, 1.85, 1.7, 0.95, CARD,
          border=BLUE, label="Learner\n(Browser)", lsz=14, lcol=DIM,
          desc="Learner browser node")
    _arrow_r(slide, 2.52, 2.18)
    _rect(slide, 2.9, 1.6, 2.1, 1.45, RGBColor(0x0D, 0x1A, 0x33),
          border=BLUE, label="React\nFrontend\n(Vite + TS)", lsz=14, lcol=BLUE,
          desc="React frontend node")
    _arrow_r(slide, 5.02, 2.18)
    _rect(slide, 5.4, 1.6, 2.2, 1.45, RGBColor(0x12, 0x10, 0x30),
          border=PURPLE, label="Express API\nBackend\n(Node.js)", lsz=14, lcol=PURPLE,
          desc="Express API backend node")

    # ── Vertical drop from Express API to horizontal bus ─────────────────
    api_cx = 6.5  # center of Express API box
    ln = slide.shapes.add_shape(
        1, Inches(api_cx - 0.01), Inches(3.05), Inches(0.02), Inches(0.3))
    ln.fill.solid(); ln.fill.fore_color.rgb = BLUE; ln.line.fill.background()
    ln.name = "Connector: Express API to bus"

    # ── Horizontal distribution bus ──────────────────────────────────────
    bus_y = 3.35
    sub_centers = [2.55, 5.0, 7.45]  # center-x of each sub-system
    bus = slide.shapes.add_shape(
        1, Inches(sub_centers[0]), Inches(bus_y),
        Inches(sub_centers[-1] - sub_centers[0]), Inches(0.03))
    bus.fill.solid(); bus.fill.fore_color.rgb = BLUE; bus.line.fill.background()
    bus.name = "Horizontal distribution bus"

    # ── Vertical drops from bus to each sub-system ───────────────────────
    for cx in sub_centers:
        ln = slide.shapes.add_shape(
            1, Inches(cx - 0.01), Inches(bus_y), Inches(0.02), Inches(0.4))
        ln.fill.solid(); ln.fill.fore_color.rgb = BLUE; ln.line.fill.background()
        ln.name = "Connector: bus to sub-system"

    # ── Tier 2: three sub-systems (evenly spaced) ───────────────────────
    subs = [
        (1.5,  "Scenario\nCorpus\n(JSON)",       BLUE,   "Scenario corpus sub-system"),
        (3.95, "Graph-RAG\n(NetworkX)",           PURPLE, "Graph-RAG sub-system"),
        (6.4,  "Ollama LLM\nCoach\n(Llama 3.2)", GREEN,  "Ollama coach sub-system"),
    ]
    for lft, label, col, desc in subs:
        _rect(slide, lft, 3.75, 2.1, 1.35, CARD, border=col,
              label=label, lsz=14, lcol=col, desc=desc)

    # ── Tier 3: evaluation + testing (spans all sub-systems) ─────────────
    _rect(slide, 1.5, 5.35, 7.0, 0.78,
          RGBColor(0x0A, 0x12, 0x28),
          border=RGBColor(0x2E, 0x44, 0x60),
          label="Evidence Gating  ·  Claim Scorer  ·  Playwright E2E (61 tests)",
          lsz=14, lcol=DIM,
          desc="Evaluation and testing layer")

    _tb(slide, 0.4, 6.35, 9.2, 0.48,
        "TypeScript  ·  React  ·  Express  ·  Node.js  ·  NetworkX  ·  Ollama  ·  Playwright",
        sz=14, color=DIM, align=PP_ALIGN.CENTER)
    _footer(slide)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — The Gap → Demo Transition (15 sec)
# ═══════════════════════════════════════════════════════════════════════════
def s05_gap_demo(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _title(slide, "Organizations Invest Heavily — But the Gap Persists")

    cards = [
        ("📄", "Documentation & Wikis\nJira · Docs",  "Hundreds of pages — useful for\nreference, not for building skill.",                  BLUE),
        ("💬", "ChatGPT · Coursera\nLinkedIn Learning", "Fast answers or courses —\nlearner skips the reasoning.",   PURPLE),
        ("🗺️", "Architecture Diagrams", "Shows the map — doesn't teach you\nto navigate under pressure.", GREEN),
    ]
    for i, (icon, heading, body, col) in enumerate(cards):
        l = 0.35 + i * 3.15
        _rect(slide, l, 1.3, 2.95, 2.8, CARD, border=col,
              desc=f"Gap card: {heading}")
        _tb(slide, l + 0.1, 1.45, 2.75, 0.65, icon,
            sz=36, align=PP_ALIGN.CENTER)
        _tb(slide, l + 0.1, 2.2, 2.75, 0.55, heading,
            sz=17, color=col, bold=True, align=PP_ALIGN.CENTER)
        _tb(slide, l + 0.1, 2.85, 2.75, 1.0, body,
            sz=15, color=DIM, align=PP_ALIGN.CENTER)

    _callout(slide, 0.6, 4.5, 8.8, 0.85,
             "The gap: no existing tool provides structured practice where learners reason through "
             "incidents, receive evidence-grounded feedback, and get coaching that nudges without revealing answers.",
             sz=16, border=PURPLE)

    _tb(slide, 1.0, 5.7, 8.0, 0.7,
        "Let me show you how.",
        sz=28, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    _tb(slide, 2.0, 6.5, 6.0, 0.4,
        "[ switching to live browser demo ]",
        sz=16, color=BLUE, align=PP_ALIGN.CENTER)
    _footer(slide)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Results + Contributions (30 sec)
# ═══════════════════════════════════════════════════════════════════════════
def s06_results(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _title(slide, "What You Just Saw — And What It Means")

    # Key metric banner
    _rect(slide, 0.4, 1.05, 9.2, 0.7, DARK2, border=GREEN,
          desc="Key metric banner")
    _tb(slide, 0.6, 1.12, 8.8, 0.55,
        "Graph-augmented retrieval + evidence gating:  0.856 → 0.963  (+12.5%)  ·  "
        "All unsupported claims eliminated",
        sz=16, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

    # ── Left column: Contributions ────────────────────────────────────────
    _rect(slide, 0.4, 1.95, 4.5, 4.8, CARD, border=BLUE,
          desc="Contributions column")
    _tb(slide, 0.6, 2.05, 4.1, 0.5, "Contributions",
        sz=20, color=BLUE, bold=True)

    contribs = [
        "Scenario-based practice framework for architecture reasoning",
        "Claim-level evidence gating (unsupported claims → 0)",
        "Guardrailed AI coaching — local LLM, 8 behavioral guardrails",
        "Graph-augmented retrieval validated by ablation study",
        "Full RQ1–RQ3 instrumentation · 61 E2E tests",
    ]
    for j, c in enumerate(contribs):
        _tb(slide, 0.6, 2.65 + j * 0.65, 4.2, 0.55, f"▸  {c}",
            sz=15, color=WHITE)

    # Scoped limitations
    _tb(slide, 0.6, 6.0, 4.2, 0.4,
        "Scoped: ablation-validated with synthetic data. No human study yet — most important next step.",
        sz=12, color=DIM, italic=True)

    # ── Right column: Ablation bar chart ──────────────────────────────────
    _rect(slide, 5.2, 1.95, 4.5, 4.8, CARD, border=PURPLE,
          desc="Ablation results column")
    _tb(slide, 5.4, 2.05, 4.1, 0.5, "Ablation Study",
        sz=20, color=PURPLE, bold=True)
    _tb(slide, 5.4, 2.55, 4.1, 0.4,
        "Mean reasoning score · 12 scenarios · 4 domains",
        sz=12, color=DIM)

    bar_data = [
        ("Flat RAG\n(baseline)",       0.856, DIM),
        ("+ Graph\ncontext",           0.921, BLUE),
        ("+ Evidence\ngating (final)", 0.963, GREEN),
    ]
    max_w = 2.5
    bar_left = 6.5
    for k, (label, score, col) in enumerate(bar_data):
        t = 3.2 + k * 1.1
        filled_w = (score / 1.0) * max_w
        _hbar(slide, bar_left, t, max_w, filled_w, 0.55, col,
              label=f"Bar: {label.replace(chr(10), ' ')}")
        _tb(slide, 5.35, t + 0.02, 1.1, 0.5, label, sz=11, color=DIM)
        _tb(slide, bar_left + filled_w + 0.06, t + 0.05, 0.7, 0.45,
            f"{score:.3f}", sz=14, color=col, bold=True)

    _footer(slide)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Future Work & Conclusion (25 sec)
# ═══════════════════════════════════════════════════════════════════════════
def s07_future(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _title(slide, "Future Work & Conclusion")

    # ── Left: Future Work ─────────────────────────────────────────────────
    _rect(slide, 0.4, 1.15, 4.5, 5.0, CARD, border=BLUE,
          desc="Future work column")
    _tb(slide, 0.6, 1.25, 4.1, 0.5, "Future Work",
        sz=20, color=BLUE, bold=True)

    fw = [
        "Controlled user study with real TPMs (ecological validity)",
        "Expanded corpus — additional domains",
        "Production retrieval (Qdrant + Neo4j)",
        "Multi-role support (SREs, new engineers)",
        "Collaborative practice sessions",
    ]
    for j, f in enumerate(fw):
        _tb(slide, 0.6, 1.9 + j * 0.82, 4.2, 0.75, f"▸  {f}",
            sz=16, color=WHITE)

    # ── Right: Conclusion ─────────────────────────────────────────────────
    _rect(slide, 5.2, 1.15, 4.5, 5.0, CARD, border=PURPLE,
          desc="Conclusion column")
    _tb(slide, 5.4, 1.25, 4.1, 0.5, "Conclusion",
        sz=20, color=PURPLE, bold=True)
    _tb(slide, 5.4, 1.9, 4.1, 4.0,
        "Architecture fluency is a skill you build\n"
        "through practice — not a documentation\n"
        "problem you solve with more wikis.\n\n"
        "OmniMentor gives new TPMs what no existing\n"
        "tool provides: a safe space to reason through\n"
        "real scenarios, fail without consequences,\n"
        "and build the mental model before day one.",
        sz=16, color=WHITE, wrap=True)

    _callout(slide, 0.4, 6.4, 9.2, 0.55,
             "Feedback welcome — especially on user study design and ecological validity.",
             border=GREEN)
    _footer(slide)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — References (5 sec)
# ═══════════════════════════════════════════════════════════════════════════
def s08_refs(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _title(slide, "References", sz=28)

    col1 = [
        "1.  Barab, S. & Squire, K. (2004). Design-Based Research:\n"
        "    Putting a stake in the ground. J. Learning Sciences, 13(1), 1–14.",
        "2.  Bandura, A. (1977). Self-efficacy: Toward a unifying theory\n"
        "    of behavioral change. Psychological Review, 84(2), 191–215.",
        "3.  Brown, J. S., Collins, A., & Duguid, P. (1989). Situated\n"
        "    cognition and the culture of learning. Ed. Researcher, 18(1), 32–42.",
        "4.  Chi, M. T. H. et al. (1989). Self-explanations: How students\n"
        "    study examples. Cognitive Science, 13(2), 145–182.",
        "5.  Collins, A., Brown, J. S., & Newman, S. E. (1989). Cognitive\n"
        "    apprenticeship. In Knowing, Learning, and Instruction. Erlbaum.",
        "6.  DBR Collective. (2003). Design-Based Research: An emerging\n"
        "    paradigm. Educational Researcher, 32(1), 5–8.",
    ]
    col2 = [
        "7.   Lave, J., & Wenger, E. (1991). Situated Learning:\n"
        "     Legitimate Peripheral Participation. Cambridge Univ. Press.",
        "8.   Paas, F. & van Merriënboer, J. J. G. (1994). Instructional\n"
        "     control of cognitive load. Ed. Psychology Review, 6(4), 351–371.",
        "9.   Sweller, J. (1988). Cognitive load during problem solving.\n"
        "     Cognitive Science, 12(2), 257–285.",
        "10.  VanLehn, K. (2011). The relative effectiveness of human\n"
        "     tutoring and ITS. Educational Psychologist, 46(4), 197–221.",
        "11.  Beyer, B. et al. (2016). Site Reliability Engineering.\n"
        "     O'Reilly Media.",
        "12.  Marteau, T. M. & Bekker, H. (1992). STAI short-form\n"
        "     development. Brit. J. Clinical Psychology, 31(3), 301–306.",
    ]
    for j, r in enumerate(col1):
        _tb(slide, 0.4, 1.1 + j * 1.0, 4.6, 0.92, r, sz=11, color=DIM)
    for j, r in enumerate(col2):
        _tb(slide, 5.1, 1.1 + j * 1.0, 4.6, 0.92, r, sz=11, color=DIM)

    _footer(slide)


# ── Narration doc generation ─────────────────────────────────────────────
def _write_speaker_notes_doc(path):
    lines = [
        "# OmniMentor Final Presentation — Speaker Notes",
        "",
        "8-slide deck. Time: slides 1-5 (~2 min) → live demo (~7 min) → slides 6-8 (~1 min).",
        "Presentation goal: seed conversation and invite feedback (GT CS6460 Goal 3).",
        "Reminder: These are practice notes, not a teleprompter. Know the story, then tell it naturally.",
        "",
    ]
    for idx, item in enumerate(SPEAKER_NOTES, start=1):
        lines.append(f"## Slide {idx}: {item['title']}")
        lines.append("")
        lines.append(item["notes"])
        lines.append("")

    with open(path, "w", encoding="utf-8") as f:
        f.write("\n".join(lines).rstrip() + "\n")


# ═══════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════
def main():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.6)

    s01_title(prs)
    s02_problem(prs)
    s03_research_design(prs)
    s04_architecture(prs)
    s05_gap_demo(prs)
    s06_results(prs)
    s07_future(prs)
    s08_refs(prs)

    for slide, note in zip(prs.slides, SPEAKER_NOTES):
        _set_notes(slide, note["notes"])

    out_dir = "personal/final-project/prep-final-video"
    os.makedirs(out_dir, exist_ok=True)

    out_pptx = f"{out_dir}/OMNIMENTOR_FINAL_PRESENTATION_V3.pptx"
    out_notes = f"{out_dir}/OMNIMENTOR_FINAL_PRESENTATION_V3_SPEAKER_NOTES.md"

    prs.save(out_pptx)
    _write_speaker_notes_doc(out_notes)
    print(f"Done — {len(prs.slides)} slides saved to {out_pptx}")
    print(f"Done — speaker notes doc saved to {out_notes}")


if __name__ == "__main__":
    main()
