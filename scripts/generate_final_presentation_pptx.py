from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def set_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide, RGBColor(0x06, 0x0B, 0x18))
    title = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    tf = title.text_frame
    tf.text = "OmniMentor"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    subtitle = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
    tf = subtitle.text_frame
    tf.word_wrap = True
    tf.text = "Guided scenario-based practice for developing architecture reasoning in new Technical Program Managers."
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(255, 255, 255)
    bottom = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
    tf = bottom.text_frame
    tf.text = "Arvind Sharma · Georgia Institute of Technology · CS 6460 · Spring 2026"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0x3B, 0x82, 0xF6)

def add_slide_2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, RGBColor(0x06, 0x0B, 0x18))
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title.text_frame.text = "Architecture Blindness"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(1))
    tf = subtitle.text_frame
    tf.word_wrap = True
    tf.text = "Nancy is a new Technical Program Manager at Omni-Mart — 30+ microservices, day-one incident meeting, no mental model yet. She faces three reinforcing gaps:"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(200, 200, 200)
    cards = [
        ("🧠", "Cognitive Overload", "Too many systems, too many names — working memory maxed out.", "Sweller, 1988"),
        ("😨", "Performance Anxiety", "Fear of asking \"obvious\" questions leads to hesitation and avoidance.", "Bandura, 1977"),
        ("👥", "Social Isolation", "In the meeting but not enough context to contribute — stuck at the periphery.", "Lave & Wenger, 1991")
    ]
    for i, (icon, heading, text, citation) in enumerate(cards):
        left = Inches(0.5 + i * 3.1)
        top = Inches(2.2)
        width = Inches(2.8)
        height = Inches(3.5)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x11, 0x18, 0x27)
        shape.line.color.rgb = RGBColor(0x3B, 0x82, 0xF6)
        tx = slide.shapes.add_textbox(left, top + Inches(0.2), width, Inches(3))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()
        p.text = icon
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(32)
        p = tf.add_paragraph()
        p.text = heading
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x8B, 0x5C, 0xF6)
        p = tf.add_paragraph()
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p = tf.add_paragraph()
        p.text = f"({citation})"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(12)
        p.font.italic = True
        p.font.color.rgb = RGBColor(150, 150, 150)
    footer = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.3))
    footer.text_frame.text = "Without structured practice, these gaps persist well beyond onboarding."
    footer.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    footer.text_frame.paragraphs[0].font.size = Pt(14)
    footer.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x3B, 0x82, 0xF6)

def add_slide_3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, RGBColor(0x06, 0x0B, 0x18))
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title.text_frame.text = "From research to design — every feature traces back"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    theory_data = [
        ("Cognitive Load Theory", "Sweller, 1988 · Paas & van Merriënboer, 1994", "Four-step practice loop that decomposes incident analysis into Brief → Investigate → Decide → Feedback"),
        ("Situated Cognition & LPP", "Brown, Collins & Duguid, 1989 · Lave & Wenger, 1991", "Realistic incident scenarios from a synthetic retail corpus (Omni-Mart, 12 scenarios, 4 domains)"),
        ("Intelligent Tutoring & Self-Explanation", "VanLehn, 2011 · Chi et al., 1989", "Claim-level evidence gating — unsupported claims flagged, reasoning gaps made visible"),
        ("Design-Based Research", "Barab & Squire, 2004 · DBR Collective, 2003", "Weekly design-evaluate-refine iterations with ablation validation")
    ]
    for i, (theory, citation, decision) in enumerate(theory_data):
        left = Inches(0.25 + i * 2.4)
        tb = slide.shapes.add_textbox(left, Inches(1.2), Inches(2.3), Inches(1.5))
        tb.text_frame.word_wrap = True
        p = tb.text_frame.paragraphs[0]
        p.text = theory
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0x8B, 0x5C, 0xF6)
        p.alignment = PP_ALIGN.CENTER
        p2 = tb.text_frame.add_paragraph()
        p2.text = citation
        p2.font.size = Pt(10)
        p2.font.italic = True
        p2.font.color.rgb = RGBColor(180, 180, 180)
        p2.alignment = PP_ALIGN.CENTER
        arrow = slide.shapes.add_textbox(left, Inches(2.6), Inches(2.3), Inches(0.4))
        arrow.text_frame.text = "↓"
        arrow.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        arrow.text_frame.paragraphs[0].font.size = Pt(24)
        arrow.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x3B, 0x82, 0xF6)
        db = slide.shapes.add_textbox(left, Inches(3.2), Inches(2.3), Inches(2.5))
        db.text_frame.word_wrap = True
        p3 = db.text_frame.paragraphs[0]
        p3.text = decision
        p3.font.size = Pt(13)
        p3.font.color.rgb = RGBColor(255, 255, 255)
        p3.alignment = PP_ALIGN.CENTER
    callout = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
    callout.text_frame.text = "+ AI coaching assistant: local LLM (Ollama / Llama 3.2) provides step-aware guidance without revealing answers."
    callout.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    callout.text_frame.paragraphs[0].font.size = Pt(14)
    callout.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x3B, 0x82, 0xF6)

def add_slide_4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, RGBColor(0x06, 0x0B, 0x18))
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title.text_frame.text = "Current tools give answers — they don't develop reasoning"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    tools = [
        ("📖", "Documentation", "Reference, not practice."),
        ("💬", "Chat / AI Assistants", "Fast answers — learner skips the reasoning."),
        ("🗺", "Architecture Diagrams", "Shows topology — doesn't teach interpretation under pressure.")
    ]
    for i, (icon, heading, text) in enumerate(tools):
        left = Inches(0.5 + i * 3.1)
        top = Inches(2.0)
        width = Inches(2.8)
        height = Inches(2.5)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x11, 0x18, 0x27)
        shape.line.color.rgb = RGBColor(0x8B, 0x5C, 0xF6)
        tx = slide.shapes.add_textbox(left, top + Inches(0.2), width, height)
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()
        p.text = icon
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(32)
        p = tf.add_paragraph()
        p.text = heading
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p = tf.add_paragraph()
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(200, 200, 200)
    callout_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(5.2), Inches(8), Inches(1.2))
    callout_box.fill.solid()
    callout_box.fill.fore_color.rgb = RGBColor(0x1F, 0x29, 0x37)
    callout_box.line.color.rgb = RGBColor(0x3B, 0x82, 0xF6)
    tf = callout_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "The gap: No existing tool provides structured practice with evidence-grounded feedback and coaching that nudges without revealing answers."
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "Let me show you what does."
    p2.font.size = Pt(18)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(0x3B, 0x82, 0xF6)
    p2.alignment = PP_ALIGN.CENTER

def add_slide_5(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, RGBColor(0x06, 0x0B, 0x18))
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    title.text_frame.text = "What you just saw — and what's next"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(0.5))
    subtitle.text_frame.text = "Key result: graph-augmented retrieval with evidence gating improved scores from 0.856 → 0.963 (+12.5%) and eliminated all unsupported claims."
    subtitle.text_frame.paragraphs[0].font.size = Pt(16)
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x3B, 0x82, 0xF6)
    left_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(4.4), Inches(4.5))
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = RGBColor(0x11, 0x18, 0x27)
    left_card.line.color.rgb = RGBColor(0x8B, 0x5C, 0xF6)
    tx_left = slide.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(4.2), Inches(4))
    tf_left = tx_left.text_frame
    tf_left.word_wrap = True
    p = tf_left.paragraphs[0]
    p.text = "Contributions"
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0x8B, 0x5C, 0xF6)
    contribs = ["Non-judgmental scenario-based practice framework", "Claim-level evidence gating", "Guardrailed AI coaching assistant (local LLM)", "Graph-augmented retrieval with ablation validation", "Full RQ1–RQ3 instrumentation · 61 E2E tests"]
    for c in contribs:
        p = tf_left.add_paragraph()
        p.text = f"• {c}"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(255, 255, 255)
    limitations = slide.shapes.add_textbox(Inches(0.5), Inches(6.1), Inches(4.4), Inches(0.5))
    limitations.text_frame.word_wrap = True
    p = limitations.text_frame.paragraphs[0]
    p.text = "Scoped Limitations: Ablation-validated, no human study yet. Synthetic-only data. Deterministic retrieval (production backends planned)."
    p.font.size = Pt(11)
    p.font.italic = True
    p.font.color.rgb = RGBColor(180, 180, 180)
    right_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.1), Inches(1.5), Inches(4.4), Inches(4.5))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = RGBColor(0x11, 0x18, 0x27)
    right_card.line.color.rgb = RGBColor(0x3B, 0x82, 0xF6)
    tx_right = slide.shapes.add_textbox(Inches(5.2), Inches(1.6), Inches(4.2), Inches(4))
    tf_right = tx_right.text_frame
    tf_right.word_wrap = True
    p = tf_right.paragraphs[0]
    p.text = "Future Work"
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0x3B, 0x82, 0xF6)
    future = ["Controlled user study with real TPMs", "Expanded scenario corpus & additional domains", "Production retrieval (Qdrant + Neo4j)", "Multi-role support (SREs, engineers)", "Collaborative practice sessions"]
    for f_work in future:
        p = tf_right.add_paragraph()
        p.text = f"• {f_work}"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(255, 255, 255)
    callout_bottom = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.4))
    callout_bottom.text_frame.text = "Architecture fluency is a learnable skill, not a documentation-access problem."
    callout_bottom.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    callout_bottom.text_frame.paragraphs[0].font.size = Pt(18)
    callout_bottom.text_frame.paragraphs[0].font.bold = True
    callout_bottom.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x8B, 0x5C, 0xF6)

def add_slide_6(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, RGBColor(0x06, 0x0B, 0x18))
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title.text_frame.text = "References"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    refs = [
        "1. Barab, S. & Squire, K. (2004). Design-Based Research: Putting a stake in the ground. Journal of the Learning Sciences, 13(1), 1–14.",
        "2. Bandura, A. (1977). Self-efficacy: Toward a unifying theory of behavioral change. Psychological Review, 84(2), 191–215.",
        "3. Brown, J. S., Collins, A., & Duguid, P. (1989). Situated cognition and the culture of learning. Educational Researcher, 18(1), 32–42.",
        "4. Chi, M. T. H. et al. (1989). Self-explanations: How students study and use examples in learning to solve problems. Cognitive Science, 13(2), 145–182.",
        "5. Collins, A., Brown, J. S., & Newman, S. E. (1989). Cognitive apprenticeship. In Knowing, Learning, and Instruction. Lawrence Erlbaum.",
        "6. The Design-Based Research Collective. (2003). Design-Based Research: An emerging paradigm. Educational Researcher, 32(1), 5–8.",
        "7. Lave, J., & Wenger, E. (1991). Situated Learning: Legitimate Peripheral Participation. Cambridge University Press.",
        "8. Paas, F. & van Merriënboer, J. J. G. (1994). Instructional control of cognitive load. Educational Psychology Review, 6(4), 351–371.",
        "9. Sweller, J. (1988). Cognitive load during problem solving. Cognitive Science, 12(2), 257–285.",
        "10. VanLehn, K. (2011). The relative effectiveness of human tutoring, intelligent tutoring systems, and other tutoring systems. Educational Psychologist, 46(4), 197–221.",
        "11. Beyer, B. et al. (2016). Site Reliability Engineering. O'Reilly Media.",
        "12. Marteau, T. M. & Bekker, H. (1992). STAI short-form development. British Journal of Clinical Psychology, 31(3), 301–306."
    ]
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(5.5))
    tf = tx.text_frame
    tf.word_wrap = True
    for r in refs:
        p = tf.add_paragraph()
        p.text = r
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(220, 220, 220)
    footer = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.4))
    footer.text_frame.text = "Arvind Sharma · CS 6460 · Georgia Institute of Technology · Spring 2026"
    footer.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    footer.text_frame.paragraphs[0].font.size = Pt(12)
    footer.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x3B, 0x82, 0xF6)

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)
add_title_slide(prs)
add_slide_2(prs)
add_slide_3(prs)
add_slide_4(prs)
add_slide_5(prs)
add_slide_6(prs)
prs.save("personal/OMNIMENTOR_FINAL_PRESENTATION_FROM_INSTRUCTIONS.pptx")
print("Presentation saved successfully.")
